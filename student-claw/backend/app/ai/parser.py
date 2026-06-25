"""
Multi-modal parsing & OCR (blueprint §5.1).

Pure extraction + normalisation. All heavy, blocking libraries (PyMuPDF/fitz,
python-pptx, Pillow) are invoked through asyncio.to_thread so the event loop is
never stalled. File download/storage lives in app.ai.storage; the caller hands
raw bytes to these extractors.

VLM pipeline
────────────
Images and scanned/image-only PDF pages are processed by Qwen 2.5 VL 72B via
the OpenRouter API (OpenAI-compatible endpoint).  The model receives a base-64
data-URI and returns structured text suitable for RAG chunking.

For PDFs, PyMuPDF text extraction is tried first (fast path). Pages that yield
no selectable text are rendered to high-resolution PNGs and sent to the VLM
(fallback path).  This covers both text-native and fully-scanned PDFs.

Returned `ParsedContent.segments` preserves natural boundaries (PDF pages /
PPTX slides) so the chunker can keep slides atomic while recursively splitting
PDFs.
"""

from __future__ import annotations

import asyncio
import base64
import logging
import re
import unicodedata
from dataclasses import dataclass, field
from typing import Optional

from tenacity import (
    retry,
    retry_if_exception_type,
    stop_after_attempt,
    wait_exponential,
)

from app.ai.config import MIME_PDF, MIME_PPTX

logger = logging.getLogger("student_claw.ai.parser")


class ParseError(Exception):
    """Raised when a document cannot be parsed (malformed / unsupported)."""


@dataclass
class ParsedContent:
    text: str                       # full normalised text
    segments: list[str] = field(default_factory=list)  # page/slide boundaries
    modality: str = "text"          # "text" | "image" | "pdf" | "pptx"
    language: Optional[str] = None


# ---------------------------------------------------------------------------
# Normalisation
# ---------------------------------------------------------------------------
_ZERO_WIDTH = dict.fromkeys(map(ord, "\u200b\u200c\u200d\ufeff\u2060"), None)
_CONTROL_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]")
_MULTI_WS_RE = re.compile(r"[ \t]{2,}")
_MULTI_NL_RE = re.compile(r"\n{3,}")
# Strip leftover Telegram/markdown entity noise that survives copy/paste.
_TG_MARKUP_RE = re.compile(r"(\*\*|__|~~|```|`)")


def normalize_text(text: Optional[str]) -> str:
    """
    Unicode-normalise (NFC), strip zero-width + control chars, remove broken
    Telegram markup, and collapse redundant whitespace. Idempotent.
    """
    if not text:
        return ""
    text = unicodedata.normalize("NFC", text)
    text = text.translate(_ZERO_WIDTH)
    text = _CONTROL_RE.sub("", text)
    text = _TG_MARKUP_RE.sub("", text)
    text = _MULTI_WS_RE.sub(" ", text)
    text = _MULTI_NL_RE.sub("\n\n", text)
    return text.strip()


def detect_language(text: str) -> Optional[str]:
    """Best-effort ISO-639-1 language detection; None if unavailable."""
    sample = text.strip()
    if len(sample) < 12:
        return None
    try:
        from langdetect import detect  # type: ignore

        return detect(sample)
    except Exception:  # pragma: no cover - optional dependency / short text
        return None


# ---------------------------------------------------------------------------
# OpenRouter / Qwen 2.5 VL client
# ---------------------------------------------------------------------------
_VLM_SYSTEM_PROMPT = (
    "You are a data extraction assistant. Extract ALL text visible in the "
    "image faithfully and completely. Preserve the logical reading order. "
    "Format tables as pipe-separated Markdown. Preserve numbered lists and "
    "bullet points. Do NOT add commentary, summaries, or any content not "
    "present in the image. Output only the extracted text."
)

_VLM_USER_PROMPT = (
    "Extract all text from this image exactly as described in your instructions."
)


def _build_vlm_client():
    """Build an async OpenAI-compatible client pointed at OpenRouter."""
    from openai import AsyncOpenAI

    from app.ai.config import get_ai_settings

    cfg = get_ai_settings()
    if not cfg.openrouter_api_key:
        raise RuntimeError(
            "OPENROUTER_API_KEY is not set. Cannot use VLM extraction."
        )
    return AsyncOpenAI(
        api_key=cfg.openrouter_api_key,
        base_url=cfg.openrouter_base_url,
    )


@retry(
    retry=retry_if_exception_type(Exception),
    stop=stop_after_attempt(3),
    wait=wait_exponential(multiplier=1, min=2, max=10),
    reraise=True,
)
async def _run_vlm_extraction(image_bytes: bytes, mime: str = "jpeg") -> str:
    """
    Send *image_bytes* to Qwen 2.5 VL via OpenRouter and return the extracted
    text string.  Retried up to 3 times with exponential back-off.

    Args:
        image_bytes: Raw image bytes (PNG, JPEG, …).
        mime: Image MIME sub-type without the "image/" prefix (default "jpeg").
              Pass "png" for PNG images.
    """
    from app.ai.config import get_ai_settings

    cfg = get_ai_settings()
    client = _build_vlm_client()

    b64 = base64.b64encode(image_bytes).decode("ascii")
    data_uri = f"data:image/{mime};base64,{b64}"

    response = await client.chat.completions.create(
        model=cfg.openrouter_model,
        messages=[
            {"role": "system", "content": _VLM_SYSTEM_PROMPT},
            {
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {"url": data_uri},
                    },
                    {"type": "text", "text": _VLM_USER_PROMPT},
                ],
            },
        ],
        max_tokens=4096,
        temperature=0.0,
    )
    return (response.choices[0].message.content or "").strip()


# ---------------------------------------------------------------------------
# Image extraction (VLM)
# ---------------------------------------------------------------------------
async def extract_text_from_image(data: bytes) -> ParsedContent:
    """
    Extract text from an image using the Qwen 2.5 VL 72B VLM via OpenRouter.
    Falls back to an empty result rather than raising so that ingestion can
    continue even if the VLM is temporarily unavailable.
    """
    # Detect PNG vs JPEG by magic bytes so we pass the right MIME to the model.
    mime = "png" if data[:4] == b"\x89PNG" else "jpeg"
    try:
        raw = await _run_vlm_extraction(data, mime=mime)
        text = normalize_text(raw)
    except Exception as exc:  # pragma: no cover – network / quota failures
        logger.error("VLM extraction failed for image (%d bytes): %s", len(data), exc)
        text = ""

    if not text:
        logger.info("VLM produced no text for image (%d bytes).", len(data))

    return ParsedContent(
        text=text,
        segments=[text] if text else [],
        modality="image",
        language=detect_language(text),
    )


# ---------------------------------------------------------------------------
# PDF (PyMuPDF / fitz) — text-native fast-path + VLM scanned-page fallback
# ---------------------------------------------------------------------------
def _pdf_pages_sync(data: bytes) -> list[str]:
    """Extract selectable text from each PDF page (fast path)."""
    import fitz  # PyMuPDF

    try:
        doc = fitz.open(stream=data, filetype="pdf")
    except Exception as exc:
        raise ParseError(f"Could not open PDF: {exc}") from exc
    pages: list[str] = []
    try:
        for page in doc:
            try:
                pages.append(page.get_text("text") or "")
            except Exception as exc:  # pragma: no cover - per-page resilience
                logger.warning("Failed to extract a PDF page: %s", exc)
                pages.append("")
    finally:
        doc.close()
    return pages


def _pdf_pages_to_pngs_sync(data: bytes, max_pages: int, zoom: float = 2.0) -> list[bytes]:
    """Render the first `max_pages` PDF pages to PNG bytes (for VLM fallback)."""
    import fitz  # PyMuPDF

    try:
        doc = fitz.open(stream=data, filetype="pdf")
    except Exception as exc:
        raise ParseError(f"Could not open PDF for rendering: {exc}") from exc
    images: list[bytes] = []
    try:
        matrix = fitz.Matrix(zoom, zoom)
        for page in doc[:max_pages]:
            pix = page.get_pixmap(matrix=matrix)
            images.append(pix.tobytes("png"))
    finally:
        doc.close()
    return images


async def render_pdf_pages_to_images(data: bytes, max_pages: int = 5) -> list[bytes]:
    """Async wrapper: render up to `max_pages` pages to PNG bytes."""
    return await asyncio.to_thread(_pdf_pages_to_pngs_sync, data, max_pages)


async def extract_text_from_pdf(data: bytes) -> ParsedContent:
    """
    Extract text from a PDF.

    Strategy:
    1. Try PyMuPDF text extraction on every page (fast, no API cost).
    2. For pages that yield no selectable text (scanned / image-only), render
       the page to a high-resolution PNG and send it to the VLM.
    3. Concatenate results, preserving page order.
    """
    raw_pages = await asyncio.to_thread(_pdf_pages_sync, data)
    pages: list[str] = []

    for page_idx, raw in enumerate(raw_pages):
        normalised = normalize_text(raw)
        if normalised:
            # Fast path: selectable text found.
            pages.append(normalised)
            continue

        # Fallback: render this page and run VLM OCR.
        logger.info(
            "PDF page %d has no selectable text — attempting VLM OCR.", page_idx + 1
        )
        try:
            png_list = await asyncio.to_thread(
                _pdf_pages_to_pngs_sync, data, max_pages=page_idx + 1, zoom=2.0
            )
            if png_list:
                page_png = png_list[-1]  # the page we just rendered
                vlm_text = await _run_vlm_extraction(page_png, mime="png")
                vlm_text = normalize_text(vlm_text)
                pages.append(vlm_text)
            else:
                pages.append("")
        except Exception as exc:
            logger.error(
                "VLM fallback failed for PDF page %d: %s", page_idx + 1, exc
            )
            pages.append("")

    pages = [p for p in pages if p]
    full = "\n\n".join(pages)
    if not full:
        logger.info("PDF contained no extractable text (%d bytes).", len(data))
    return ParsedContent(
        text=full, segments=pages, modality="pdf", language=detect_language(full)
    )


# ---------------------------------------------------------------------------
# PPTX (python-pptx) — slides are semantically atomic
# ---------------------------------------------------------------------------
def _pptx_slides_sync(data: bytes) -> list[str]:
    from io import BytesIO

    from pptx import Presentation

    try:
        prs = Presentation(BytesIO(data))
    except Exception as exc:
        raise ParseError(f"Could not open PPTX: {exc}") from exc

    slides: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line:
                        parts.append(line)
            if shape.has_table:  # type: ignore[attr-defined]
                for row in shape.table.rows:
                    cells = [c.text for c in row.cells if c.text]
                    if cells:
                        parts.append(" | ".join(cells))
        slides.append("\n".join(parts))
    return slides


async def extract_text_from_pptx(data: bytes) -> ParsedContent:
    raw_slides = await asyncio.to_thread(_pptx_slides_sync, data)
    slides = [normalize_text(s) for s in raw_slides]
    slides = [s for s in slides if s]
    full = "\n\n".join(slides)
    return ParsedContent(
        text=full, segments=slides, modality="pptx", language=detect_language(full)
    )


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------
async def parse_document(data: bytes, mime_type: Optional[str], filename: str = "") -> ParsedContent:
    """
    Route a document to the right extractor by MIME type, with a filename-based
    fallback. Raises ParseError for unsupported types.
    """
    mime = (mime_type or "").lower()
    name = (filename or "").lower()

    if mime == MIME_PDF or name.endswith(".pdf"):
        return await extract_text_from_pdf(data)
    if mime == MIME_PPTX or name.endswith(".pptx"):
        return await extract_text_from_pptx(data)

    raise ParseError(f"Unsupported document type: mime={mime_type!r} name={filename!r}")
